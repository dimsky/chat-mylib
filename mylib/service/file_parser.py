# -*- coding: utf-8 -*-

import os
import PyPDF2
import docx
from mylib.service.exception import BaseException, ExceededMaxPagesError, ExceededMaxWordsError, UnsupportError
import openpyxl
import hashlib
import unicodedata
from pptx import Presentation

def to_halfwidth(text):
    # 将字符串中的所有字符转换为半角
    return unicodedata.normalize('NFKC', text)
                                 

class FileStruct:
    file_path = ''
    type = ''
    pages = []
    full_text = ''
    full_text_md5 = ''

    def __init__(self, file_path, type, pages, full_text,full_text_md5):
        self.file_path = file_path
        self.type = type
        self.pages = pages
        self.full_text = full_text
        self.full_text_md5 = full_text_md5


class FileParser:
    SUPPORT_DOC_MAX_WORDS = 500000
    SUPPORT_DOC_MAX_PAGES = 500
    # SUPPORT_EXTS = ['.txt', '.csv', 'md', '.pdf', '.doc', '.docx', '.xlsx', 'xls']
    #word 没办法分页读取,暂时不支持
    SUPPORT_EXTS = ['.txt', '.csv', '.md', '.pdf','.xlsx', '.ppt', '.pptx']

    def __get_parser_by_ext(self, ext):
        if ext.lower() in ['.txt', '.md', '.csv']:
            return TXTParser()
        elif ext.lower() in ['.pdf']:
            return PDFParser()
        elif ext.lower() in ['.doc', '.docx']:
            return WordParser()
        elif ext.lower() in ['.xlsx']:
            return ExcelParser()
        elif ext.lower() in ['.ppt', '.pptx']:
            return PptParser()

    # 将文件内容解析为文本
    def parse(self, file_path)->FileStruct:
        # 根据扩展名生成策略类
        ext = os.path.splitext(file_path)[1].lower()
        if ext not in self.SUPPORT_EXTS:
            raise UnsupportError('不支持的文件格式！ 目前支持的文件格式有:' + ','.join(self.SUPPORT_EXTS))

        paraser = self.__get_parser_by_ext(ext)
        return paraser.parse(file_path)


class TXTParser(FileParser):
    def parse(self, file_path)->FileStruct:
        with open(file_path, 'r') as file:
            full_text = file.read()
            if (len(full_text) > self.SUPPORT_DOC_MAX_WORDS):
                # 抛出异常
                raise ExceededMaxWordsError('超出最大字数:' + str(self.SUPPORT_DOC_MAX_WORDS))
            md5 = hashlib.md5(full_text.encode('utf-8')).hexdigest()
        return FileStruct(file_path=file_path, type='txt', pages=[full_text], full_text=full_text,full_text_md5 = md5)

class PDFParser(FileParser):
    def parse(self, file_path)->FileStruct:
        full_text = ''
        pages = []
        # 打开PDF文件
        with open(file_path, 'rb') as file:
            # 创建PdfFileReader对象
            pdf = PyPDF2.PdfReader(file_path)
            # 获取PDF文件的页面数量
            num_pages = len(pdf.pages)
            if num_pages > self.SUPPORT_DOC_MAX_PAGES:
                raise ExceededMaxPagesError('超出最大页数:' + str(self.SUPPORT_DOC_MAX_PAGES))

            # 读取每一页的内容
            for page in range(num_pages):
                page_text = pdf.pages[page].extract_text()
                page_text = to_halfwidth(page_text)
                pages.append(page_text)
                full_text += page_text

            if (len(full_text) > self.SUPPORT_DOC_MAX_WORDS):
                raise ExceededMaxPagesError('超出最大字数:' + str(self.SUPPORT_DOC_MAX_WORDS))
            md5 = hashlib.md5(full_text.encode('utf-8')).hexdigest()
            # 返回解析后的文本内容
            return FileStruct(file_path=file_path, type='pdf', pages=pages, full_text=full_text, full_text_md5 = md5)


class WordParser(FileParser):
    def parse(self, file_path)->FileStruct:
        full_text = ''
        # 打开Word文件
        doc = docx.Document(file_path)
        # 读取每一段的内容
        for para in doc.paragraphs:
            full_text += para.text
        if (len(full_text) > self.SUPPORT_DOC_MAX_WORDS):
            raise ExceededMaxWordsError('超出最大字数:' + str(self.SUPPORT_DOC_MAX_WORDS))
        md5 = hashlib.md5(full_text.encode('utf-8')).hexdigest()
        # 返回解析后的文本内容
        return FileStruct(file_path=file_path, type='word', pages=[full_text], full_text=full_text, full_text_md5 = md5)


class ExcelParser(FileParser):
    def parse(self, file_path)->FileStruct:
        full_text = ''
        # 打开Excel文件
        wb = openpyxl.load_workbook(file_path)
        # 选择第一个工作表
        sheet = wb.active
        # 读取每一个单元格的内容
        for row in sheet.rows:
            row_str = '|'
            for cell in row:
                row_str += str(cell.value) + '|'
            full_text += row_str + '\n'
        if (len(full_text) > self.SUPPORT_DOC_MAX_WORDS):
            raise ExceededMaxWordsError('超出最大字数:' + str(self.SUPPORT_DOC_MAX_WORDS))
        md5 = hashlib.md5(full_text.encode('utf-8')).hexdigest()
        # 返回解析后的文本内容
        return FileStruct(file_path=file_path, type='excel', pages=[full_text], full_text=full_text, full_text_md5 = md5)
    
class PptParser(FileParser):
    def parse(self, file_path)->FileStruct:
        full_text = ''
        pages = []
        prs = Presentation(file_path)
        num_pages = len(prs.slides)
        for page in range(num_pages):
            slide = prs.slides[page]
            page_text = ""
             # 遍历幻灯片上的每个形状
            for shape in slide.shapes:
                # 如果形状是文本框
                if shape.has_text_frame:
                    # 遍历文本框的每一段
                    for paragraph in shape.text_frame.paragraphs:
                        # 遍历段落中的每个文本片段
                        for run in paragraph.runs:
                            # 输出文本片段的内容
                            page_text += run.text + " "
            pages.append(page_text)
            full_text += page_text
        if (len(full_text) > self.SUPPORT_DOC_MAX_WORDS):
            raise ExceededMaxWordsError('超出最大字数:' + str(self.SUPPORT_DOC_MAX_WORDS))
        md5 = hashlib.md5(full_text.encode('utf-8')).hexdigest()
        # 返回解析后的文本内容
        return FileStruct(file_path=file_path, type='ppt', pages=pages, full_text=full_text, full_text_md5 = md5)


        